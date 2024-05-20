import pandas as pd
#import snowflake.connector
import streamlit as st
from streamlit_dynamic_filters import DynamicFilters
from st_aggrid import AgGrid
from st_aggrid.grid_options_builder import GridOptionsBuilder
from st_aggrid import GridUpdateMode, DataReturnMode
import warnings
from yahooquery import Ticker
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
from datetime import date
from PIL import Image
import requests
import os
from io import BytesIO
from langchain.chat_models import ChatOpenAI
from langchain.schema import HumanMessage, SystemMessage
import traceback
import re
import ast

# hide future warnings (caused by st_aggrid)
warnings.simplefilter(action='ignore', category=FutureWarning)

# set page layout and define basic variables
st.set_page_config(layout="wide", page_icon='⚡', page_title="Instant Insight")
path = os.path.dirname(__file__)
today = date.today()


def resize_image(url):
    """function to resize logos while keeping aspect ratio. Accepts URL as an argument and return an image object"""

    # Open the image file
    image = Image.open(requests.get(url, stream=True).raw)

    # if a logo is too high or too wide then make the background container twice as big
    if image.height > 140:
        container_width = 220 * 2
        container_height = 140 * 2

    elif image.width > 220:
        container_width = 220 * 2
        container_height = 140 * 2
    else:
        container_width = 220
        container_height = 140

    # Create a new image with the same aspect ratio as the original image
    new_image = Image.new('RGBA', (container_width, container_height))

    # Calculate the position to paste the image so that it is centered
    x = (container_width - image.width) // 2
    y = (container_height - image.height) // 2

    # Paste the image onto the new image
    new_image.paste(image, (x, y))
    return new_image


def add_image(slide, image, left, top, width):
    """function to add an image to the PowerPoint slide and specify its position and width"""
    slide.shapes.add_picture(image, left=left, top=top, width=width)


def replace_text(replacements, slide):
    """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text


def get_stock(ticker, period, interval):
    """function to get stock data from Yahoo Finance. Takes ticker, period and interval as arguments and returns a DataFrame"""
    hist = ticker.history(period=period, interval=interval)
    hist = hist.reset_index()
    # capitalize column names
    hist.columns = [x.capitalize() for x in hist.columns]
    return hist


def plot_graph(df, x, y, title, name):
    """function to plot a line graph. Takes DataFrame, x and y axis, title and name as arguments and returns a Plotly figure"""
    fig = px.line(df, x=x, y=y, template='simple_white',
                  title='<b>{} {}</b>'.format(name, title))
    fig.update_traces(line_color='#A27D4F')
    fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
    return fig


def peers_plot(df, name, metric):
    """function to plot a bar chart with peers. Takes DataFrame, name, metric and ticker as arguments and returns a Plotly figure"""

    # drop rows with missing metrics
    df.dropna(subset=[metric], inplace=True)

    df_sorted = df.sort_values(metric, ascending=False)

    # iterate over the labels and add the colors to the color mapping dictionary, hightlight the selected ticker
    color_map = {}
    for label in df_sorted['Company Name']:
        if label == name:
            color_map[label] = '#A27D4F'
        else:
            color_map[label] = '#D9D9D9'

    fig = px.bar(df_sorted, y='Company Name', x=metric, template='simple_white', color='Company Name',
                 color_discrete_map=color_map,
                 orientation='h',
                 title='<b>{} {} vs Peers FY22</b>'.format(name, metric))
    fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)', showlegend=False, yaxis_title='')
    return fig


def esg_plot(name, df):
    # Define colors for types
    colors = {name: '#A27D4F', 'Peer Group': '#D9D9D9'}

    # Creating the bar chart
    fig = go.Figure()
    for type in df['Type'].unique():
        fig.add_trace(go.Bar(
            x=df[df['Type'] == type]['variable'],
            y=df[df['Type'] == type]['value'],
            name=type,
            text=df[df['Type'] == type]['value'],
            textposition='outside',
            marker_color=colors[type]
        ))
    fig.update_layout(
        height=700,
        width=1000,
        barmode='group',
        title="ESG Score vs Peers Average",
        xaxis_title="",
        yaxis_title="Score",
        legend_title="Type",
        xaxis=dict(tickangle=0),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)')
    return fig


def get_financials(df, col_name, metric_name):
    """function to get financial metrics from a DataFrame. Takes DataFrame, column name and metric name as arguments and returns a DataFrame"""
    metric = df.loc[:, ['asOfDate', col_name]]
    metric_df = pd.DataFrame(metric).reset_index()
    metric_df.columns = ['Symbol', 'Year', metric_name]

    return metric_df


def generate_gpt_response(gpt_input, max_tokens, api_key, llm_model):
    """function to generate a response from GPT-3. Takes input and max tokens as arguments and returns a response"""
    # Create an instance of the OpenAI class
    chat = ChatOpenAI(openai_api_key=api_key, model=llm_model,
                      temperature=0, max_tokens=max_tokens)

    # Generate a response from the model
    response = chat.predict_messages(
        [SystemMessage(content='You are a helpful expert in finance, market and company research.'
                               'You also have exceptional skills in selling B2B software products.'),
         HumanMessage(
             content=gpt_input)])

    return response.content.strip()


def dict_from_string(response):
    """function to parse GPT response with competitors tickers and convert it to a dict"""
    # Find a substring that starts with '{' and ends with '}', across multiple lines
    match = re.search(r'\{.*?\}', response, re.DOTALL)

    dictionary = None
    if match:
        try:
            # Try to convert substring to dict
            dictionary = ast.literal_eval(match.group())
        except (ValueError, SyntaxError):
            # Not a dictionary
            return None
    return dictionary


def extract_comp_financials(tkr, comp_name, dict):
    """function to extract financial metrics for competitors. Takes a ticker as an argument and appends financial metrics to dict"""
    ticker = Ticker(tkr)
    income_df = ticker.income_statement(frequency='a', trailing=False)

    subset = income_df.loc[:, ['asOfDate', 'TotalRevenue', 'SellingGeneralAndAdministration']].reset_index()

    # keep only 2022 data
    subset = subset[subset['asOfDate'].dt.year == 2022].sort_values(by='asOfDate', ascending=False).head(1)

    # get values
    total_revenue = subset['TotalRevenue'].values[0]
    sg_and_a = subset['SellingGeneralAndAdministration'].values[0]

    # calculate sg&a as a percentage of total revenue
    sg_and_a_pct = round(sg_and_a / total_revenue * 100, 2)

    # add values to dictionary
    dict[comp_name]['Total Revenue'] = total_revenue
    dict[comp_name]['SG&A % Of Revenue'] = sg_and_a_pct


def convert_to_nested_dict(input_dict, nested_key):
    """function to convert a dictionary to a nested dictionary. Takes a dictionary and a nested key as arguments and returns a dictionary"""
    output_dict = {}
    for key, value in input_dict.items():
        output_dict[key] = {nested_key: value}
    return output_dict


def shorten_summary(text):
    # Split the text into sentences using a regular expression pattern
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)

    # Return the first two sentences or less if there are fewer sentences
    sen = ' '.join(sentences[:2])

    # if the summary is less than 350 characters, return the summary
    if len(sen) <= 400:
        return sen
    else:
        # if the summary is more than 350 characters, return the first 350 characters and truncate the last word
        truncated_sen = text[:400].rsplit(' ', 1)[0] + '...'
        return truncated_sen


def peers_summary(df, selected_ticker):
    df = df[df['Ticker'] != selected_ticker]

    for tkr in df['Ticker']:
        try:
            profile = Ticker(tkr).asset_profile
            summary = profile[tkr]['longBusinessSummary']
            website = profile[tkr]['website']

            # keep only the first two sentences of the summary
            short_summary = shorten_summary(summary)
            logo_url = 'https://logo.clearbit.com/' + website

            # append short summary and logo_url to the df
            df.loc[df['Ticker'] == tkr, 'Summary'] = short_summary
            df.loc[df['Ticker'] == tkr, 'Logo'] = logo_url
        except:
            continue
    # drop rows with missing summary
    df = df.dropna(subset=['Summary'])
    return df


def fix_text_capitalization(text):
    fixed_text = text.lower().capitalize()
    return fixed_text


def replace_multiple_symbols(string):
    """function to fix description from yahoo finance, sometimes it has multiple dots at the end of the string"""
    string = string.replace(':', '')
    pattern = r'\.{2,}$'  # Matches two or more consecutive dots at the end of the string
    replacement = '.'

    # Check if the string ends with multiple symbols
    if re.search(pattern, string):
        # Replace multiple symbols with a single symbol
        string = re.sub(pattern, replacement, string)
    return string


def no_data_plot():
    """plot to return when there is no data available"""
    # Create a blank figure with a transparent background
    fig = go.Figure()

    # Add a text annotation for "NO DATA AVAILABLE" at the center of the plot
    fig.add_annotation(
        x=0.5,
        y=0.5,
        xref='paper',
        yref='paper',
        text='NO DATA AVAILABLE',
        showarrow=False,
        font=dict(size=26, color='black'),
    )

    # Customize layout to have a transparent background
    fig.update_layout(
        paper_bgcolor='rgba(0,0,0,0)',  # Transparent background
        plot_bgcolor='rgba(0,0,0,0)',  # Transparent plot area
        xaxis_showgrid=False,  # Hide x-axis gridlines
        yaxis_showgrid=False,  # Hide y-axis gridlines
        xaxis=dict(visible=False),  # Hide x-axis labels and ticks
        yaxis=dict(visible=False),  # Hide y-axis labels and ticks
    )
    return fig


#conn = st.connection("snowflake")
#df = conn.query("SELECT * from prospects LIMIT 1000;", ttl=600)
df = pd.read_csv('prospects.csv')

# Fix column names. Replace underscore with space, lowercase column names, and capitalize first words
df.columns = df.columns.str.replace('_', ' ').str.lower().str.title()

with st.sidebar:
    openai_key = st.text_input(label="Your OpenAI API key", help="Your API key is not stored anywhere")
    llm_model = st.selectbox(label="Choose a model", options=["gpt-3.5-turbo", "gpt-4-turbo", "gpt-4", "gpt-4o"])

# create sidebar filters
st.sidebar.write('**Use filters to select prospects** 👇')

# display dynamic multi select filters
dynamic_filters = DynamicFilters(df, filters=['Sector', 'Industry', 'Prospect Status', 'Product'])
dynamic_filters.display_filters(location='sidebar')
df_filtered = dynamic_filters.filter_df()


with st.sidebar:
    st.markdown('''The dataset is taken from [Kaggle](https://www.kaggle.com/datasets/aramacus/usa-public-companies) and slightly modified for the purpose of this app.
    ''', unsafe_allow_html=True)
    st.markdown('''[GitHub Repo](https://github.com/arsentievalex/instant-insight-web-app)''', unsafe_allow_html=True)
    st.markdown('''The app created by [Oleksandr Arsentiev](https://twitter.com/alexarsentiev) for the purpose of
    Streamlit Summit Hackathon''', unsafe_allow_html=True)

##############################################################################################################

st.title('Welcome to the Instant Insight App!⚡')

with st.expander('What is this app about?'):
    st.write('''
    This app is designed to generate an instant company research.\n
    In a matter of few clicks, a user gets a PowerPoint presentation with the company overview, SWOT analysis, financials, and value propostion tailored for the selling product. 
    The app works with the US public companies.

    Use Case Example:\n
    Imagine working in sales for a B2B SaaS company that has hundreds of prospects and offers the following products: 
    Accounting and Planning Software, CRM, Chatbot, and Cloud Data Storage.
    You are tasked to do a basic prospect research and create presentations for your team. The prospects data is stored in a Snowflake database that feeds your CRM system.
    You can use this app to quickly filter the prospects by sector, industry, prospect status, and product. 
    Next, you can select the prospect you want to include in the presentation and click the button to generate the presentation.
    And...that's it! You have the slides ready to be shared with your team.

    Tech Stack:\n
    • Database - Snowflake via Snowflake Connector\n
    • Data Processing - Pandas\n
    • Research Data - Yahoo Finance via Yahooquery, GPT 3.5 via LangChain\n
    • Visualization - Plotly\n
    • Frontend - Streamlit, AgGrid\n
    • Presentation - Python-pptx\n
    ''')

num_of_cust = df_filtered.shape[0]
st.metric(label='Number of Prospects', value=num_of_cust)

# button to create slides
ui_container = st.container()
with ui_container:
    submit = st.button(label='Generate Presentation')

# select columns to show
df_filtered = df_filtered[['Company Name', 'Sector', 'Industry', 'Prospect Status', 'Product']]

# creating AgGrid dynamic table and setting configurations
gb = GridOptionsBuilder.from_dataframe(df_filtered)
gb.configure_selection(selection_mode="single", use_checkbox=True)
gb.configure_column(field='Company Name', width=270)
gb.configure_column(field='Sector', width=260)
gb.configure_column(field='Industry', width=350)
gb.configure_column(field='Prospect Status', width=270)
gb.configure_column(field='Product', width=240)

gridOptions = gb.build()

response = AgGrid(
    df_filtered,
    gridOptions=gridOptions,
    enable_enterprise_modules=False,
    height=600,
    update_mode=GridUpdateMode.SELECTION_CHANGED,
    data_return_mode=DataReturnMode.FILTERED_AND_SORTED,
    fit_columns_on_grid_load=False,
    theme='alpine',
    allow_unsafe_jscode=True
)

# get selected rows
response_df = pd.DataFrame(response["selected_rows"])

# if user input is empty and button is clicked then show warning
if submit and response_df.empty:
    with ui_container:
        st.warning("Please select a prospect!")

elif submit and openai_key == "":
    with ui_container:
        st.warning("Please input your OpenAI API key")

# if user input is not empty and button is clicked then generate slides
elif submit and response_df is not None:
    with ui_container:
        with st.spinner('Generating awesome slides for you...⏳'):

            try:
                # define variables for selected prospect
                company_name = response_df['Company Name'].values[0]
                product = response_df['Product'].values[0]

                # join df with response_df to get a ticker of selected prospect
                df_ticker = pd.merge(df, response_df, left_on='Company Name', right_on='Company Name')
                selected_ticker = df_ticker['Ticker'].values[0]

                # open presentation template
                pptx = path + '//' + 'template.pptx'
                prs = Presentation(pptx)

                ticker = Ticker(selected_ticker)

                # get stock info
                name = ticker.price[selected_ticker]['shortName']
                sector = ticker.summary_profile[selected_ticker]['sector']
                industry = ticker.summary_profile[selected_ticker]['industry']
                employees = ticker.summary_profile[selected_ticker]['fullTimeEmployees']
                country = ticker.summary_profile[selected_ticker]['country']
                city = ticker.summary_profile[selected_ticker]['city']
                website = ticker.summary_profile[selected_ticker]['website']
                summary = ticker.summary_profile[selected_ticker]['longBusinessSummary']
                logo_url = 'https://logo.clearbit.com/' + website

                # declare pptx variables
                title_slide = prs.slides[0]
                summary_slide = prs.slides[1]
                s_w_slide = prs.slides[2]
                vp_slide = prs.slides[4]
                key_people_slide = prs.slides[5]
                comp_slide = prs.slides[6]
                esg_slide = prs.slides[7]

                # initiate a dictionary of placeholders and values to replace
                replaces_1 = {
                    '{company}': name,
                    '{date}': today}

                replaces_2 = {
                    '{c}': name,
                    '{s}': sector,
                    '{i}': industry,
                    '{co}': country,
                    '{ci}': city,
                    '{ee}': "{:,}".format(employees),
                    '{w}': website,
                    '{summary}': summary
                }

                # run the function to replace placeholders with values
                replace_text(replaces_1, title_slide)
                replace_text(replaces_2, summary_slide)

                # check if a logo ulr returns code 200 (working link)
                if requests.get(logo_url).status_code == 200:
                    # create logo image object
                    logo = resize_image(logo_url)
                    logo.save('logo.png')
                    logo_im = 'logo.png'

                    # add logo to the slide
                    add_image(prs.slides[1], image=logo_im, left=Inches(1.2), width=Inches(2), top=Inches(0.5))
                    os.remove('logo.png')

                ##############################################################################################################
                # create slides with financial plots
                # get financial data
                fin_df = ticker.all_financial_data()

                # plot stock price
                stock_df = get_stock(ticker=ticker, period='5y', interval='1mo')
                stock_fig = plot_graph(df=stock_df, x='Date', y='Open', title='Stock Price USD', name=name)

                stock_fig.write_image("stock.png")
                stock_im = 'stock.png'

                add_image(prs.slides[3], image=stock_im, left=Inches(1.8), width=Inches(4.5), top=Inches(0.5))
                os.remove('stock.png')

                # plot revenue
                rev_df = get_financials(df=fin_df, col_name='TotalRevenue', metric_name='Total Revenue')
                rev_fig = plot_graph(df=rev_df, x='Year', y='Total Revenue', title='Total Revenue USD', name=name)

                rev_fig.write_image("rev.png")
                rev_im = 'rev.png'

                add_image(prs.slides[3], image=rev_im, left=Inches(1.8), width=Inches(4.5), top=Inches(3.8))
                os.remove('rev.png')

                # plot market cap
                debt_df = get_financials(df=fin_df, col_name='TotalDebt', metric_name='Total Debt')
                debt_fig = plot_graph(df=debt_df, x='Year', y='Total Debt', title='Total Debt USD', name=name)

                debt_fig.write_image("marketcap.png")
                debt_im = 'marketcap.png'

                add_image(prs.slides[3], image=debt_im, left=Inches(7.3), width=Inches(4.5), top=Inches(0.5))
                os.remove('marketcap.png')

                # plot ebitda
                # adding try and except because some companies like banks don't have EBITDA data
                try:
                    ebitda_df = get_financials(df=fin_df, col_name='NormalizedEBITDA', metric_name='EBITDA')
                    ebitda_fig = plot_graph(df=ebitda_df, x='Year', y='EBITDA', title='EBITDA USD', name=name)

                    ebitda_fig.write_image("ebitda.png")
                    ebitda_im = 'ebitda.png'

                    add_image(prs.slides[3], image=ebitda_im, left=Inches(7.3), width=Inches(4.5), top=Inches(3.8))
                    os.remove('ebitda.png')
                except:
                    pass

                ############################################################################################################
                
                # create competitors slide
                input_competitors = """What are the top competitors of {} company with ticker {}?
                                    Provide up to 4 most relevant public US competitors comparable by revenue and market cap.
                                    Return output as a Python dictionary with company name as key and ticker as value.
                                    Do not return anything else."""

                # format template with company name and ticker
                input_competitors = input_competitors.format(name, selected_ticker)

                # return response from GPT-3
                gpt_comp_response = generate_gpt_response(gpt_input=input_competitors, max_tokens=250, api_key=openai_key, llm_model=llm_model)

                # extract dictionary from response
                peers_dict = dict_from_string(gpt_comp_response)

                # check if any competitors were returned
                if peers_dict is not None:

                    # convert dict to nested dict to later hold financial data
                    peers_dict_nested = convert_to_nested_dict(input_dict=peers_dict, nested_key='Ticker')

                    # add current ticker to the list
                    peers_dict_nested[name] = {'Ticker': selected_ticker}

                    # extract financial data for each competitor
                    for key, value in peers_dict_nested.items():
                        try:
                            extract_comp_financials(tkr=value['Ticker'], comp_name=key, dict=peers_dict_nested)
                        # if ticker is not found, drop it from the peers dict and continue
                        except:
                            del peers_dict[key]
                            continue

                    # create a dataframe with peers financial data
                    peers_df = pd.DataFrame.from_dict(peers_dict_nested, orient='index')
                    peers_df = peers_df.reset_index().rename(columns={'index': 'Company Name'})

                    # plot revenue vs peers graph
                    sg_and_a_peers_fig = peers_plot(df=peers_df, name=name, metric='SG&A % Of Revenue')

                    sg_and_a_peers_fig.write_image("sg_and_a_peers.png")
                    sg_and_a_peers_im = 'sg_and_a_peers.png'

                    add_image(prs.slides[6], image=sg_and_a_peers_im, left=Inches(0.8), width=Inches(4.8),
                              top=Inches(0.5))
                    os.remove('sg_and_a_peers.png')

                    # plot operating expenses vs peers graph
                    rev_peers_fig = peers_plot(df=peers_df, name=name, metric='Total Revenue')
                    rev_peers_fig.write_image("rev_peers.png")
                    rev_peers_im = 'rev_peers.png'

                    add_image(prs.slides[6], image=rev_peers_im, left=Inches(0.8), width=Inches(4.8), top=Inches(3.8))
                    os.remove('rev_peers.png')

                    # get competitor company descriptions
                    peers_summary_df = peers_summary(df=peers_df, selected_ticker=selected_ticker)

                    # create a list of competitor descriptions and logos
                    summary_list = peers_summary_df['Summary'].tolist()
                    logo_list = peers_summary_df['Logo'].tolist()

                    # if there are less than 4 competitors, add empty strings to the list
                    if len(summary_list) < 4:
                        summary_list += [''] * (4 - len(summary_list))

                    # unpack list into variables
                    c1, c2, c3, c4 = summary_list

                    # initiate a dictionary of placeholders and values to replace
                    replaces_5 = {
                        '{a}': c1,
                        '{b}': c2,
                        '{c}': c3,
                        '{d}': c4}

                    # replace placeholders with values
                    replace_text(replaces_5, comp_slide)

                    top_row = Inches(0.7)

                    for l in logo_list:
                        # check if a logo ulr returns code 200 (working link)
                        if requests.get(l).status_code == 200:
                            # create logo image object
                            logo = resize_image(l)
                            logo.save('logo.png')
                            logo_im = 'logo.png'

                            # add logo to the slide
                            add_image(comp_slide, image=logo_im, left=Inches(5.4), width=Inches(1.2), top=top_row)
                            top_row += Inches(1.8)
                            os.remove('logo.png')

                ############################################################################################################

                # create strengths and weaknesses slide

                input_swot = f"""
                Create a SWOT analysis of {name} company with ticker {selected_ticker}
                Return output as a Python dictionary with the following keys: Strengths, Weaknesses, 
                Opportunities, Threats. The values should be a brief description of each in string format.
                Do not return anything else. Be concise and specific.
                """
                
                # return response from GPT-3
                gpt_swot = generate_gpt_response(gpt_input=input_swot, max_tokens=1000, api_key=openai_key, llm_model=llm_model)

                # extract dictionary from response
                swot_dict = dict_from_string(gpt_swot)

                # create title for the slide
                swot_title = 'SWOT Analysis of {}'.format(name)

                # initiate a dictionary of placeholders and values to replace
                replaces_3 = {
                    '{s}': swot_dict['Strengths'],
                    '{w}': swot_dict['Weaknesses'],
                    '{o}': swot_dict['Opportunities'],
                    '{t}': swot_dict['Threats'],
                    '{swot_title}': swot_title}

                # run the function to replace placeholders with values
                replace_text(replaces_3, s_w_slide)

                ############################################################################################
                
                # create value prop slide
                input_vp = """"Create a brief value proposition using Value Proposition Canvas framework for {product} for 
                {name} company with ticker {ticker} that operates in {industry} industry.
                Return output as a Python dictionary with the following keys: Pains, Gains, Gain Creators, 
                Pain Relievers as a keys and text as values. Be specific and concise. Do not return anything else."""

                input_vp = input_vp.format(product=product, name=name, ticker=selected_ticker, industry=industry)

                # return response from GPT-3
                gpt_value_prop = generate_gpt_response(gpt_input=input_vp, max_tokens=1000, api_key=openai_key, llm_model=llm_model)

                # extract dictionary from response
                value_prop_dict = dict_from_string(gpt_value_prop)

                vp_title = 'Value Proposition of {} for {}'.format(product, name)

                # initiate a dictionary of placeholders and values to replace
                replaces_4 = {
                    '{p}': value_prop_dict['Pains'],
                    '{g}': value_prop_dict['Gains'],
                    '{gc}': value_prop_dict['Gain Creators'],
                    '{pr}': value_prop_dict['Pain Relievers'],
                    '{vp_title}': vp_title}

                # run the function to replace placeholders with values
                replace_text(replaces_4, vp_slide)

                ############################################################################################

                # key people slide
                key_people = ticker.asset_profile[selected_ticker]['companyOfficers']

                # create title, name and age lists from key_people
                kp_titles = []
                kp_names = []
                kp_age = []

                for i in range(4):
                    try:
                        kp_titles.append(key_people[i]['title'])
                        kp_names.append(key_people[i]['name'])
                        kp_age.append(key_people[i]['age'])
                    except:
                        kp_titles.append('')
                        kp_names.append('')
                        kp_age.append('')

                replaces_6 = {
                    '{t1}': kp_titles[0],
                    '{t2}': kp_titles[1],
                    '{t3}': kp_titles[2],
                    '{t4}': kp_titles[3],
                    '{n1}': kp_names[0],
                    '{n2}': kp_names[1],
                    '{n3}': kp_names[2],
                    '{n4}': kp_names[3],
                    '{a1}': kp_age[0],
                    '{a2}': kp_age[1],
                    '{a3}': kp_age[2],
                    '{a4}': kp_age[3],
                    '{company_name}': name}

                # run the function to replace placeholders with values
                replace_text(replaces_6, key_people_slide)

                ############################################################################################
            
                # corporate news
                news_df = ticker.corporate_events

                # sort by date descending
                news_df = news_df.sort_values(by=['date'], ascending=False)
                # reset index
                news_df.reset_index(inplace=True)

                # keep only top three rows
                news_df = news_df.head(3)

                # clean description column
                news_df['fixed_description'] = news_df['description'].apply(fix_text_capitalization).apply(
                    replace_multiple_symbols).apply(shorten_summary)

                # Remove the timestamp and keep only the date
                news_df['date'] = news_df['date'].dt.date

                # drop duplicates in headline columns
                news_df.drop_duplicates(subset=['headline'], inplace=True)

                # drop rows with empty headline
                news_df.dropna(subset=['headline'], inplace=True)

                # create title, name and age lists from key_people
                news_headlines = []
                news_dates = []
                news_desc = []

                for i in range(3):
                    try:
                        news_headlines.append(news_df['headline'][i])
                        news_dates.append(news_df['date'][i])
                        news_desc.append(news_df['fixed_description'][i])
                    except:
                        news_headlines.append('')
                        news_dates.append('')
                        news_desc.append('')

                replaces_7 = {
                    '{h1}': news_headlines[0],
                    '{h2}': news_headlines[1],
                    '{h3}': news_headlines[2],
                    '{d1}': news_dates[0],
                    '{d2}': news_dates[1],
                    '{d3}': news_dates[2],
                    '{desc1}': news_desc[0],
                    '{desc2}': news_desc[1],
                    '{desc3}': news_desc[2]}

                # run the function to replace placeholders with values
                replace_text(replaces_7, key_people_slide)

                ############################################################################################
            
                # ESG slide
                # get ESG scores from Yahoo Finance
                try:
                    esg_scores = ticker.esg_scores

                    # get esg score data for the company
                    total_esg = esg_scores[selected_ticker]['totalEsg']
                    governance_score = esg_scores[selected_ticker]['governanceScore']
                    environment_score = esg_scores[selected_ticker]['environmentScore']
                    social_score = esg_scores[selected_ticker]['socialScore']

                    # get peer group data
                    peer_group = esg_scores[selected_ticker]['peerGroup']
                    peers_total_esg_avg = esg_scores[selected_ticker]['peerEsgScorePerformance']['avg']
                    peers_governance_avg = esg_scores[selected_ticker]['peerGovernancePerformance']['avg']
                    peers_env_avg = esg_scores[selected_ticker]['peerEnvironmentPerformance']['avg']
                    peers_social_avg = esg_scores[selected_ticker]['peerSocialPerformance']['avg']

                    esg_dict = {'Type': [name, 'Peer Group'],
                                'Total ESG Score': [round(total_esg, 2), round(peers_total_esg_avg, 2)],
                                'Governance Score': [round(governance_score, 2), round(peers_governance_avg, 2)],
                                'Environment Score': [round(environment_score, 2), round(peers_env_avg, 2)],
                                'Social Score': [round(social_score, 2), round(peers_social_avg, 2)]}

                    esg_df = pd.DataFrame(esg_dict)
                    # Pivot DataFrame
                    esg_df_melted = esg_df.melt(id_vars='Type',
                                                value_vars=['Total ESG Score', 'Governance Score', 'Environment Score',
                                                            'Social Score'])
                    # run function to generate a bar chart with esg comparison
                    esg_fig = esg_plot(name=name, df=esg_df_melted)

                    esg_fig.write_image("esg.png")
                    esg_im = 'esg.png'

                    add_image(esg_slide, image=esg_im, left=Inches(3), width=Inches(7.7), top=Inches(1.6))
                    os.remove('esg.png')

                except Exception as e:
                    # if the ESG data is not avaialble, return plot with a message (preferably the slide would be deleted, but there is no option for that in python-pptx)
                    no_data_fig = no_data_plot()
                    no_data_fig.write_image("no_data.png")
                    no_data_im = 'no_data.png'

                    add_image(esg_slide, image=no_data_im, left=Inches(3), width=Inches(7.7), top=Inches(1.6))
                    os.remove('no_data.png')

                # replace title with company name
                replaces_8 = {'{company_name}': name}
                # run the function to replace placeholders with values
                replace_text(replaces_8, esg_slide)

                ###########################################################################################
            
                # create file name
                filename = '{} {}.pptx'.format(name, today)

                # save presentation as binary output
                binary_output = BytesIO()
                prs.save(binary_output)

                # display success message and download button
                with ui_container:
                    st.success('The slides have been generated! :tada:')

                    st.download_button(label='Click to download PowerPoint',
                                       data=binary_output.getvalue(),
                                       file_name=filename)

            # if there is any error, display an error message
            except Exception as e:
                with ui_container:
                    st.write(e)
                    # get more details on error
                    st.write(traceback.format_exc())
                    st.error("Oops, something went wrong, please try again or select a different prospect.")
